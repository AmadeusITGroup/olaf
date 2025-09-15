import os
import re
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

def parse_presentation_file(file_path):
    """Parse a markdown presentation file and extract slide data."""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    # Extract title
    title_match = re.search(r'^# (.+?)$', content, re.MULTILINE)
    presentation_title = title_match.group(1) if title_match else "Presentation"
    
    # Extract slides
    slides = []
    slide_pattern = r'### Slide \d+: (.+?)\n\*\*Layout\*\*: (.+?)\n\*\*Content\*\*:\n(.*?)(?=\n\*Image Prompt|\n###|\Z)'
    
    for match in re.finditer(slide_pattern, content, re.DOTALL):
        slide_title = match.group(1)
        layout = match.group(2)
        content_text = match.group(3).strip()
        
        # Extract image prompt if exists
        image_prompt = None
        image_match = re.search(r'\*Image Prompt[^:]*: (.+?)(?=\n###|\Z)', content[match.end():], re.DOTALL)
        if image_match:
            image_prompt = image_match.group(1).strip()
        
        slides.append({
            'title': slide_title,
            'layout': layout,
            'content': content_text,
            'image_prompt': image_prompt
        })
    
    return presentation_title, slides

def create_presentation_from_file(input_file_path, output_dir="findings_dir/presentations"):
    """Create a PowerPoint presentation from a parsed markdown file."""
    
    # Parse the input file
    presentation_title, slides = parse_presentation_file(input_file_path)
    
    # Create a presentation object
    prs = Presentation()
    
    for i, slide_data in enumerate(slides):
        # Choose layout based on slide type
        if slide_data['layout'] == 'Title Slide':
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            
            title.text = slide_data['title']
            if subtitle and slide_data['content']:
                # Extract subtitle from content
                lines = slide_data['content'].split('\n')
                subtitle_text = next((line.replace('- Subtitle: ', '') for line in lines if 'Subtitle:' in line), '')
                subtitle.text = subtitle_text
                
        else:
            # Content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = slide_data['title']
            content.text = slide_data['content']
        
        # Add image prompt as a comment/note if it exists
        if slide_data['image_prompt']:
            # Add image prompt as slide notes
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = f"Image Prompt: {slide_data['image_prompt']}"
    
    # Generate output filename
    safe_title = re.sub(r'[^\w\s-]', '', presentation_title).strip()
    safe_title = re.sub(r'[-\s]+', '_', safe_title)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    output_filename = f"{safe_title}_{timestamp}.pptx"
    
    # Ensure output directory exists
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, output_filename)
    
    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved as: {output_path}")
    return output_path

def main():
    """Main function to handle command line usage."""
    import sys
    
    if len(sys.argv) < 2:
        print("Usage: python generate_dynamic_pptx.py <input_markdown_file> [output_directory]")
        print("Example: python generate_dynamic_pptx.py findings_dir/my_presentation.md")
        return
    
    input_file = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else "findings_dir/presentations"
    
    if not os.path.exists(input_file):
        print(f"Error: Input file '{input_file}' not found.")
        return
    
    try:
        output_path = create_presentation_from_file(input_file, output_dir)
        print(f"Successfully created presentation: {output_path}")
    except Exception as e:
        print(f"Error creating presentation: {str(e)}")

if __name__ == "__main__":
    main()